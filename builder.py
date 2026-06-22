import io
import requests
from pptx import Presentation
from pptx.util import Inches

class PptxBuilder:
    def __init__(self, template_path=None):
        # テンプレートが指定されていればそれを、なければ白紙のデフォルトファイルを作成
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.current_slide = None
        self.body_shape = None

    def _add_new_slide(self, title_text: str):
        """新しいスライドを追加し、タイトルを設定"""
        # 「タイトルとコンテンツ」のレイアウトを使用
        slide_layout = self.prs.slide_layouts[1]
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        
        # タイトル枠にテキストを挿入
        title_shape = self.current_slide.shapes.title
        if title_shape:
            title_shape.text = title_text
        
        # コンテンツ枠を取得
        self.body_shape = self.current_slide.placeholders[1]

    def build(self, parsed_data: list, output_file: str):
        """解析済みデータを読み込み、スライドを組み立てるメイン処理"""
        
        for item in parsed_data:
            image_type = item["type"]

            # 1. 新しいスライドの作成（見出し1, 2をトリガーにする）
            if image_type in ["heading_1", "heading_2"]:
                self._add_new_slide(item["text"])

            # 2. テキストの追加（段落や箇条書き）
            elif image_type in ["paragraph", "bulleted_list_item", "numbered_list_item"]:
                # スライドが1枚もない状態でテキストが来た場合のフェイルセーフ
                if not self.current_slide:
                    self._add_new_slide("無題のスライド")
                
                text_frame = self.body_shape.text_frame
                p = text_frame.add_paragraph()
                p.text = item["text"]
                
                # 箇条書きは少しインデントを下げる
                if image_type in ["bulleted_list_item", "numbered_list_item"]:
                    p.level = 1
                else:
                    p.level = 0

            # 3. 画像の追加
            elif image_type == "image":
                if not self.current_slide:
                    self._add_new_slide("画像スライド")
                
                # 画像をダウンロードして貼り付け
                self._add_image_to_slide(item["url"])

        # 最後にファイルを保存
        self.prs.save(output_file)

    def _add_image_to_slide(self, image_url: str):
        """画像をダウンロードしてスライドに挿入する"""
        try:
            # URLから画像データを取得
            response = requests.get(image_url)
            response.raise_for_status() # エラーがあればここでストップ
            
            # メモリ上に一時保存
            image_stream = io.BytesIO(response.content)
            
            # スライドの左から1インチ、上から2インチの位置に、幅4インチで挿入
            self.current_slide.shapes.add_picture(
                image_stream, 
                Inches(1), Inches(2), width=Inches(4)
            )
        except Exception as e:
            print(f"画像の挿入に失敗しました: {e}")