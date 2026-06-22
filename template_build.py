"""
template_builder.py
python main.py を実行する前に一度だけ実行し、template.pptx を生成するスクリプト。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os

# ============================================================
# 定数定義
# ============================================================

# スライドサイズ（4:3）
SLIDE_WIDTH  = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

# カラー
NAVY      = RGBColor(0x1F, 0x4E, 0x79)
BLUE_LINE = RGBColor(0x2E, 0x75, 0xB6)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GRAY      = RGBColor(0x80, 0x80, 0x80)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# フォント
FONT_JP = "MS Gothic"
FONT_EN = "Arial"

# 出力先
OUTPUT_PATH = "assets/template.pptx"


# ============================================================
# ユーティリティ関数
# ============================================================

def set_font(run, size_pt, bold=False, color=None, font_en=FONT_EN, font_jp=FONT_JP):
    """runのフォントを設定する"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = font_en
    # 日本語フォントの設定
    rPr = run._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_jp)


def add_textbox(slide, left, top, width, height, text, size_pt,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                font_en=FONT_EN, font_jp=FONT_JP):
    """スライドにテキストボックスを追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold=bold, color=color, font_en=font_en, font_jp=font_jp)
    return txBox


def add_line(slide, x1, y1, x2, y2, color, width_pt=1.5):
    """スライドに直線を追加する"""
    from pptx.util import Pt as PtUtil
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """スライドに矩形を追加する"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ============================================================
# 各レイアウトの装飾要素をスライドに追加する関数
# ============================================================

def decorate_title(slide):
    """
    レイアウト0：タイトルスライド
    - 白背景にネイビーの長方形を約35度回転させて左下に配置
    - 日付テキストボックス
    - タイトルテキストボックス
    """
    W = int(SLIDE_WIDTH)
    H = int(SLIDE_HEIGHT)

    # 回転角度（60000分の1度単位 / 1度 = 60000）
    # 35度回転
    ROT = 35 * 60000

    # 長方形のサイズ（スライドより大きめにして回転後も隙間が出ないようにする）
    rect_cx = int(W * 1.3)
    rect_cy = int(H * 1.0)

    # 配置中心：スライド左下寄り
    rect_x = int(W * -0.2)
    rect_y = int(H * 0.15)

    rect_xml = f"""
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="100" name="NavyRect"/>
        <p:cNvSpPr><a:spLocks noChangeArrowheads="1"/></p:cNvSpPr>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm rot="{ROT}">
          <a:off x="{rect_x}" y="{rect_y}"/>
          <a:ext cx="{rect_cx}" cy="{rect_cy}"/>
        </a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill>
          <a:srgbClr val="1F4E79"/>
        </a:solidFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr/>
        <a:lstStyle/>
        <a:p/>
      </p:txBody>
    </p:sp>
    """
    rect_elem = etree.fromstring(rect_xml.strip())
    slide.shapes._spTree.append(rect_elem)

    # 日付テキストボックス（長方形の上に重ねる）
    add_textbox(
        slide,
        left=Inches(1.0), top=Inches(3.8),
        width=Inches(5), height=Inches(0.5),
        text="YYYY/MM/DD",
        size_pt=18, bold=False, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # タイトルテキストボックス
    add_textbox(
        slide,
        left=Inches(1.0), top=Inches(4.5),
        width=Inches(5), height=Inches(0.8),
        text="タイトル",
        size_pt=28, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )


def decorate_toc_overview(slide):
    """
    レイアウト1：全体目次スライド
    - 「目次」タイトル
    - タイトル下の水色アンダーライン
    - 大題目リスト用プレースホルダー
    """
    W = SLIDE_WIDTH

    # 「目次」タイトル
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(0.3),
        width=Inches(9), height=Inches(0.7),
        text="目次",
        size_pt=28, bold=True, color=BLACK
    )

    # 水色アンダーライン
    add_line(
        slide,
        x1=Inches(0.5), y1=Inches(1.1),
        x2=Inches(9.5), y2=Inches(1.1),
        color=BLUE_LINE, width_pt=2.0
    )

    # 大題目リスト用テキストボックス（プレースホルダー代わり）
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(1.3),
        width=Inches(9), height=Inches(5.8),
        text="（目次リストがここに入ります）",
        size_pt=20, bold=False, color=GRAY
    )


def decorate_toc_chapter(slide):
    """
    レイアウト2：大題目ごとの目次スライド
    - 「目次」タイトル
    - 水色アンダーライン
    - 動的コンテンツ用の空エリア（builder.pyが書き込む）
    """
    # 「目次」タイトル
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(0.3),
        width=Inches(9), height=Inches(0.7),
        text="目次",
        size_pt=28, bold=True, color=BLACK
    )

    # 水色アンダーライン
    add_line(
        slide,
        x1=Inches(0.5), y1=Inches(1.1),
        x2=Inches(9.5), y2=Inches(1.1),
        color=BLUE_LINE, width_pt=2.0
    )
    # コンテンツエリアはbuilder.pyが動的に生成するため、ここでは何も置かない


def decorate_content(slide, section_count: int):
    """
    レイアウト3〜5：本編スライド
    - 左上に大題目（小・グレー）
    - スライドタイトル
    - 水色アンダーライン
    - section_count個の小セクション見出し＋本文エリア
    """
    W = SLIDE_WIDTH
    H = SLIDE_HEIGHT

    # 大題目（左上・小・グレー）
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(0.1),
        width=Inches(9), height=Inches(0.3),
        text="大題目",
        size_pt=11, bold=False, color=GRAY
    )

    # スライドタイトル
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(0.4),
        width=Inches(9), height=Inches(0.6),
        text="小題目タイトル",
        size_pt=24, bold=True, color=BLACK
    )

    # 水色アンダーライン
    add_line(
        slide,
        x1=Inches(0.5), y1=Inches(1.05),
        x2=Inches(9.5), y2=Inches(1.05),
        color=BLUE_LINE, width_pt=1.5
    )

    # セクション数に応じてコンテンツエリアを分割
    content_top    = Inches(1.2)
    content_bottom = Inches(7.3)
    content_height = content_bottom - content_top
    section_height = content_height / section_count
    section_gap    = Inches(0.15)

    for i in range(section_count):
        top = content_top + section_height * i

        # グレーのサイドバー（小セクション見出しの左端）
        add_rect(
            slide,
            left=Inches(0.5), top=top,
            width=Inches(0.08), height=Inches(0.4),
            fill_color=BLUE_LINE
        )

        # 小セクション見出し
        add_textbox(
            slide,
            left=Inches(0.7), top=top,
            width=Inches(8.8), height=Inches(0.4),
            text=f"小セクション{i + 1}",
            size_pt=14, bold=True, color=BLACK
        )

        # 本文エリア
        add_textbox(
            slide,
            left=Inches(0.7), top=top + Inches(0.45),
            width=Inches(8.8), height=section_height - Inches(0.45) - section_gap,
            text="本文テキスト",
            size_pt=12, bold=False, color=BLACK
        )


# ============================================================
# メイン処理
# ============================================================

def build_template():
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 既存のレイアウトをすべて使用して6種類のスライドを生成
    # blank レイアウト（インデックス6）をベースに使う
    blank_layout = prs.slide_layouts[6]

    layout_funcs = [
        ("TITLE",        decorate_title),
        ("TOC_OVERVIEW", decorate_toc_overview),
        ("TOC_CHAPTER",  decorate_toc_chapter),
        ("CONTENT_1",    lambda s: decorate_content(s, 1)),
        ("CONTENT_2",    lambda s: decorate_content(s, 2)),
        ("CONTENT_3",    lambda s: decorate_content(s, 3)),
    ]

    for name, func in layout_funcs:
        slide = prs.slides.add_slide(blank_layout)
        func(slide)
        print(f"  ✓ {name} を生成しました")

    os.makedirs("output", exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"\n✅ {OUTPUT_PATH} を生成しました（スライド数: {len(prs.slides)}）")


if __name__ == "__main__":
    print("テンプレートを生成しています...")
    build_template()
